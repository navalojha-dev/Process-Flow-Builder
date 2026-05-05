#!/usr/bin/env python3
"""
render.py — Shipsy process-flow PPT generator (v2).

Usage:
    python3 render.py <flow.json> <output.pptx>

Reads a flow definition (see schema.json) and produces a 3-slide PPTX:
  1. Process Flow with Shipsy   — programmatic, dynamic per use case
                                   (rows, steps, icons all flow-driven)
  2. Process Mapping & Capability Enhancement matrix
  3. AI Use Cases · AgentFleet for {client}

Slide 1 no longer uses a fixed PPTX template. The agent decides how
many rows, how many steps per row, and which icon represents each step
based on the inferred client profile and use case.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
ICONS_DIR = HERE / "icons"
ICON_MANIFEST = ICONS_DIR / "manifest.json"

# Shipsy palette (sampled from Skynet deck + matrix template screenshot)
SHIPSY_NAVY = RGBColor(0x1F, 0x3A, 0x5F)
SHIPSY_BLUE = RGBColor(0x1F, 0x4E, 0x79)
SHIPSY_CYAN = RGBColor(0x4F, 0xB0, 0xC6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x55, 0x55, 0x55)
# Matrix-slide column header colours (matched to user screenshot)
MATRIX_PROC_BG = RGBColor(0x2A, 0x2A, 0x2A)   # near-black "Business Process"
MATRIX_ASIS_BG = RGBColor(0xA9, 0x37, 0x29)   # brick red "As Is"
MATRIX_TOBE_BG = RGBColor(0x2C, 0x49, 0x6A)   # navy "To-Be"
MATRIX_IMPACT_BG = RGBColor(0x59, 0x9C, 0x38) # forest "Impact"
# Soft row backgrounds
MATRIX_ROW_BG = RGBColor(0xF6, 0xF6, 0xF6)
MATRIX_PROC_COL_BG = RGBColor(0xEE, 0xEE, 0xEE)
# AI card phase tag colors
PHASE_NOW = RGBColor(0x59, 0x9C, 0x38)        # green
PHASE_PHASE2 = RGBColor(0xE0, 0x95, 0x1B)     # amber
PHASE_FUTURE = RGBColor(0x77, 0x77, 0x77)     # gray


# ---------------------------------------------------------------------------
# Shape walking + text substitution
# ---------------------------------------------------------------------------

def walk_shapes(shapes):
    """Yield every shape including those nested inside groups."""
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from walk_shapes(shape.shapes)
        else:
            yield shape


def replace_run_text(shape, new_text: str) -> bool:
    """Replace the text in a shape preserving the first run's formatting.

    Returns True if substitution occurred. Strategy: keep first paragraph +
    first run, set its text to ``new_text``, drop all other runs/paragraphs.
    """
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    if not tf.paragraphs:
        return False
    p0 = tf.paragraphs[0]
    if not p0.runs:
        # No runs — fall back to .text setter (loses formatting but works)
        tf.text = new_text
        return True
    p0.runs[0].text = new_text
    # Remove additional runs from first paragraph
    for run in list(p0.runs[1:]):
        run._r.getparent().remove(run._r)
    # Remove all subsequent paragraphs
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    return True


# ---------------------------------------------------------------------------
# Icon resolution: name (or alias) -> path
# ---------------------------------------------------------------------------
_ICON_MANIFEST_CACHE = None


def _load_manifest():
    global _ICON_MANIFEST_CACHE
    if _ICON_MANIFEST_CACHE is None:
        try:
            _ICON_MANIFEST_CACHE = json.loads(ICON_MANIFEST.read_text())
        except Exception:
            _ICON_MANIFEST_CACHE = {"icons": {}, "aliases": {}}
    return _ICON_MANIFEST_CACHE


def _resolve_icon(name):
    """Resolve an icon name (or alias) to a file path. Returns None if unknown."""
    if not name:
        return None
    direct = ICONS_DIR / f"{name}.png"
    if direct.exists():
        return direct
    manifest = _load_manifest()
    alias = manifest.get("aliases", {}).get(name)
    if alias:
        aliased = ICONS_DIR / f"{alias}.png"
        if aliased.exists():
            return aliased
    return None


# ---------------------------------------------------------------------------
# Wrapper slide builders (Learnings + Impact & Summary)
# ---------------------------------------------------------------------------

def _set_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _add_textbox(slide, left, top, width, height, text, *, font_size=14,
                 bold=False, color=DARK_TEXT, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_title_bar(slide, slide_w, title_text, subtitle_text=None):
    bar_h = Inches(1.0) if subtitle_text else Inches(0.85)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, bar_h)
    _set_solid(bar, SHIPSY_NAVY)
    _add_textbox(slide, Inches(0.4), Inches(0.15), slide_w - Inches(0.8),
                 Inches(0.45), title_text, font_size=24, bold=True, color=WHITE)
    if subtitle_text:
        # Truncate aggressively — subtitle is for context only and the
        # full info lives in the snapshot box on slide 2.
        if len(subtitle_text) > 95:
            subtitle_text = subtitle_text[:92] + "…"
        _add_textbox(slide, Inches(0.4), Inches(0.58), slide_w - Inches(0.8),
                     Inches(0.4), subtitle_text, font_size=11,
                     color=RGBColor(0xCC, 0xDD, 0xEE))


def _estimate_wrapped_lines(text: str, width_in: float, font_pt: int) -> int:
    """Heuristic line-count estimator. Body fonts at 11-13pt have an average
    glyph width ≈ 0.0072 * font_pt inches. Slightly conservative so we err
    on giving a bullet more space rather than less."""
    if width_in <= 0:
        return 1
    char_w = max(0.005, 0.0072 * font_pt)
    chars_per_line = max(1, int(width_in / char_w))
    # Word-aware: bump line count when a single long word wouldn't fit but
    # we'd otherwise allocate exactly on a boundary.
    n_chars = len(text)
    return max(1, -(-n_chars // chars_per_line))


def _add_bullet_list(slide, items, left, top, width, height,
                     *, bullet_color=SHIPSY_BLUE, font_size=13):
    """Variable-height bullet list. Estimates each bullet's wrap count and
    places them sequentially. If total estimated height exceeds the
    available space, font and per-bullet height scale down proportionally."""
    n = len(items)
    if n == 0:
        return
    if n >= 7:
        font_size = min(font_size, 11)
    elif n >= 6:
        font_size = min(font_size, 12)

    # Inches values for layout math
    text_w_in = max(0.5, (width - Inches(0.35)) / 914400.0)
    avail_h_in = max(0.5, height / 914400.0)
    line_h_in = font_size * 0.0185 + 0.04  # px-per-pt approximation
    pad_in = 0.10  # gap between bullets
    min_h_in = line_h_in + pad_in

    heights = []
    for txt in items:
        lines = _estimate_wrapped_lines(txt, text_w_in, font_size)
        heights.append(max(min_h_in, lines * line_h_in + pad_in))

    total_required = sum(heights)
    if total_required > avail_h_in:
        scale = avail_h_in / total_required
        heights = [h * scale for h in heights]
        # If we had to compress, bump font down one notch
        font_size = max(10, font_size - 1)

    y_cursor_emu = top
    for txt, h_in in zip(items, heights):
        h_emu = Inches(h_in)
        # Bullet dot — vertically centred on the first line
        dot_d = Inches(0.12)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + Inches(0.05),
            y_cursor_emu + Inches(0.13),
            dot_d, dot_d,
        )
        _set_solid(dot, bullet_color)
        _add_textbox(slide, left + Inches(0.30), y_cursor_emu,
                     width - Inches(0.35), h_emu,
                     txt, font_size=font_size, color=DARK_TEXT)
        y_cursor_emu += h_emu


def _set_cell_fill(cell, rgb):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb


def _set_cell_text(cell, text, *, font_size=11, bold=False, color=DARK_TEXT,
                   align_left=True, vertical_center=True):
    """Replace the cell's text with a single run while preserving the cell's
    own fill. Returns the text frame for further customisation."""
    cell.text = ""  # clear default
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align_left:
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    if vertical_center:
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return tf


def _set_cell_bullets(cell, items, *, font_size=9, color=DARK_TEXT,
                      bullet_char="• "):
    """Render a list of strings as paragraph bullets inside a table cell."""
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    # Tight margins for table cells
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.06)
    cell.margin_bottom = Inches(0.06)

    # Drop autofit for first paragraph
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(2)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = bullet_char + item
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        r.font.bold = False


def render_process_flow_v2(prs, flow):
    """Slide 1 — programmatic, dynamic Process Flow with Shipsy.

    No fixed PPTX template. The agent decides:
      - how many rows (4-7)
      - what each row is called
      - how many steps per row (2-6)
      - which icon represents each step (with fallback to labeled box)
      - day-marker tag per row (optional)
      - bottom strip (3 labels)
    """
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    sw, sh = prs.slide_width, prs.slide_height

    pf = flow.get("process_flow") or {}
    title = pf.get("title", "Process Flow with Shipsy")
    rows = pf.get("rows", []) or []

    if not rows:
        # Defensive fallback so we don't render an empty slide.
        rows = [{
            "name": "(no rows in flow JSON)",
            "steps": [{"label": "—", "icon": ""}],
        }]

    # ----- 1. Title (Skynet-style: navy text on white background) -----
    title_h = Inches(0.65)
    _add_textbox(slide, Inches(0.30), Inches(0.10), Inches(7.5), Inches(0.55),
                 title, font_size=28, bold=True, color=SHIPSY_NAVY)

    # Shipsy mark in the top-right (chevron + "Shipsy" wordmark)
    shipsy_mark = ICONS_DIR / "shipsy-logo-mark.png"
    if shipsy_mark.exists():
        slide.shapes.add_picture(str(shipsy_mark),
                                 Inches(8.85), Inches(0.18),
                                 height=Inches(0.40))
    else:
        chev = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      Inches(8.65), Inches(0.22),
                                      Inches(0.30), Inches(0.24))
        _set_solid(chev, SHIPSY_BLUE)
        _add_textbox(slide, Inches(9.00), Inches(0.18),
                     Inches(1.00), Inches(0.32),
                     "Shipsy", font_size=16, bold=True, color=SHIPSY_NAVY)

    # ----- 2. Body -----
    body_top = title_h + Inches(0.08)
    body_bottom = sh - Inches(0.15)  # no bottom strip — full body to bottom
    body_h_total = body_bottom - body_top

    sidebar_w = Inches(1.55)
    body_x = sidebar_w
    body_w = sw - sidebar_w - Inches(0.20)  # right margin only

    n_rows = len(rows)
    row_h = body_h_total / n_rows

    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # Adapt fonts to row count
    if n_rows >= 7:
        sb_font, label_font = 9, 8
        icon_frac = 0.42
    elif n_rows == 6:
        sb_font, label_font = 10, 8
        icon_frac = 0.45
    elif n_rows == 5:
        sb_font, label_font = 11, 9
        icon_frac = 0.48
    else:
        sb_font, label_font = 12, 10
        icon_frac = 0.55

    # Skynet-style graduated sidebar palette: darkest navy at the top,
    # ending at a saturated medium blue (NOT pale grey-white) so the
    # sidebar holds its colour presence all the way down.
    def _sidebar_color(i, n):
        if n <= 1:
            return RGBColor(0x1F, 0x3A, 0x5F)
        t = i / (n - 1)
        r1, g1, b1 = 0x1F, 0x3A, 0x5F   # deep navy (top)
        r2, g2, b2 = 0x6E, 0x8C, 0xB8   # medium blue (bottom) — still saturated
        return RGBColor(int(r1 + (r2 - r1) * t),
                        int(g1 + (g2 - g1) * t),
                        int(b1 + (b2 - b1) * t))

    for i, row in enumerate(rows):
        y = body_top + row_h * i

        # ---- Sidebar (row label) — Skynet-style graduated colour band -----
        sb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, sidebar_w, row_h)
        _set_solid(sb, _sidebar_color(i, n_rows))
        sb.line.fill.background()  # no border — colour bands are the divider

        # Dashed horizontal divider in the body, separating rows
        if i > 0:
            div = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                sidebar_w + Inches(0.05), y,
                sw - Inches(0.20), y,
            )
            div.line.color.rgb = RGBColor(0xD9, 0xDF, 0xE6)
            div.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            div.line.width = Pt(0.5)

        sb_tb = slide.shapes.add_textbox(Inches(0.05), y + Inches(0.03),
                                         sidebar_w - Inches(0.10),
                                         row_h - Inches(0.06))
        sb_tf = sb_tb.text_frame
        sb_tf.word_wrap = True
        sb_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        sb_tf.margin_left = Inches(0.04)
        sb_tf.margin_right = Inches(0.04)
        sp = sb_tf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = row.get("name", "")
        sr.font.size = Pt(sb_font)
        sr.font.bold = True
        sr.font.color.rgb = WHITE

        # ---- Steps ----
        steps = row.get("steps") or []
        n_steps = max(1, len(steps))
        step_w = body_w / n_steps

        # Icon size: fits in a cell, capped
        icon_size = min(row_h * icon_frac, step_w * 0.45, Inches(0.55))
        icon_y = y + Inches(0.05)

        prev_icon_right_x = None
        prev_icon_cy = None

        for j, step in enumerate(steps):
            sx = body_x + step_w * j
            cx = sx + step_w / 2
            icon_path = _resolve_icon(step.get("icon"))

            if icon_path:
                ix = cx - icon_size / 2
                slide.shapes.add_picture(str(icon_path), ix, icon_y,
                                         width=icon_size, height=icon_size)
                icon_left_x = ix
                icon_right_x = ix + icon_size
                icon_cy = icon_y + icon_size / 2
            else:
                # Fallback: rounded rectangle with first 2 chars of label
                box_w = icon_size * 1.3
                box_x = cx - box_w / 2
                fbox = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    box_x, icon_y, box_w, icon_size,
                )
                _set_solid(fbox, SHIPSY_CYAN)
                fbox.line.color.rgb = SHIPSY_NAVY
                icon_left_x = box_x
                icon_right_x = box_x + box_w
                icon_cy = icon_y + icon_size / 2

            # Step label below
            label_y = icon_y + icon_size + Inches(0.04)
            label_h = max(Inches(0.20), y + row_h - label_y - Inches(0.02))
            lb = slide.shapes.add_textbox(sx + Inches(0.03), label_y,
                                          step_w - Inches(0.06), label_h)
            ltf = lb.text_frame
            ltf.word_wrap = True
            ltf.vertical_anchor = MSO_ANCHOR.TOP
            ltf.margin_top = Inches(0.0)
            ltf.margin_left = Inches(0.02)
            ltf.margin_right = Inches(0.02)
            lp = ltf.paragraphs[0]
            lp.alignment = PP_ALIGN.CENTER
            lr = lp.add_run()
            lr.text = step.get("label", "")
            lr.font.size = Pt(label_font)
            lr.font.color.rgb = DARK_TEXT

            # Skynet-style connector: a single solid block-arrow spanning
            # the gap between consecutive icons (thin body + triangle head).
            if prev_icon_right_x is not None:
                gap_start = prev_icon_right_x + Inches(0.05)
                gap_end = icon_left_x - Inches(0.02)
                if gap_end - gap_start > Inches(0.12):
                    arrow_color = RGBColor(0x9E, 0xAE, 0xC2)
                    arrow_h = Inches(0.16)
                    arrow = slide.shapes.add_shape(
                        MSO_SHAPE.RIGHT_ARROW,
                        gap_start, icon_cy - arrow_h / 2,
                        gap_end - gap_start, arrow_h,
                    )
                    _set_solid(arrow, arrow_color)
                    arrow.line.fill.background()

            prev_icon_right_x = icon_right_x
            prev_icon_cy = icon_cy

    return slide


def add_process_mapping_slide(prs, flow):
    """Slide: Process Mapping & Capability Enhancement with Shipsy.

    Layout follows the user-provided template:
      - Title bar with Shipsy mark in top-right
      - 4-column table: Business Process | As Is | To-Be | Impact
      - Coloured header row + 5-7 data rows
    """
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    sw, sh = prs.slide_width, prs.slide_height

    # --- Title bar (white background, navy text) ---
    title_h = Inches(0.70)
    # Title sized so it never reaches the right-side Shipsy mark
    _add_textbox(slide, Inches(0.35), Inches(0.20),
                 Inches(8.05), Inches(0.46),
                 "Process Mapping & Capability Enhancement with Shipsy",
                 font_size=18, bold=True, color=SHIPSY_NAVY)

    # Top-right compact Shipsy mark — chevron + small wordmark on one line
    chev = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  Inches(8.55), Inches(0.28),
                                  Inches(0.30), Inches(0.28))
    _set_solid(chev, SHIPSY_BLUE)
    _add_textbox(slide, Inches(8.90), Inches(0.26),
                 Inches(1.05), Inches(0.32),
                 "Shipsy", font_size=16, bold=True, color=SHIPSY_NAVY)

    # Thin accent line under the title
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.35), title_h,
                                    Inches(9.30), Inches(0.02))
    _set_solid(accent, RGBColor(0xCC, 0xCC, 0xCC))

    # --- Build the matrix as a fixed-grid (manual rectangles) ---
    # python-pptx native tables auto-grow rows when content overflows; we
    # need strict heights here so all rows fit in the slide. So we draw
    # each cell as a rectangle + textbox at computed coordinates.
    rows_data = flow.get("process_mapping", [])
    if not rows_data:
        rows_data = [{
            "process": "(no process_mapping in flow JSON)",
            "as_is": ["—"], "to_be": ["—"], "impact": ["—"],
        }]
    n_data = len(rows_data)

    # Geometry
    grid_left = Inches(0.30)
    grid_top = title_h + Inches(0.18)
    grid_w = sw - Inches(0.60)
    grid_h = sh - grid_top - Inches(0.20)

    col_w_proc = Inches(1.55)
    rest_w = grid_w - col_w_proc
    col_w_data = rest_w / 3.0
    col_xs = [grid_left,
              grid_left + col_w_proc,
              grid_left + col_w_proc + col_w_data,
              grid_left + col_w_proc + 2 * col_w_data]
    col_ws = [col_w_proc, col_w_data, col_w_data, col_w_data]

    header_h = Inches(0.42)
    body_avail = grid_h - header_h
    row_h = body_avail / max(1, n_data)

    # Pick body font + bullet cap that will actually fit each row
    if n_data <= 4:
        body_font = 11
        max_bullets = 5
    elif n_data == 5:
        body_font = 10
        max_bullets = 4
    elif n_data == 6:
        body_font = 9
        max_bullets = 4
    else:
        body_font = 8
        max_bullets = 3

    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    def _cell_rect(left, top, width, height, fill_rgb):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        _set_solid(rect, fill_rgb)
        rect.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        rect.line.width = Pt(0.5)
        return rect

    def _header_text(left, top, width, height, label):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = WHITE

    def _proc_text(left, top, width, height, label):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(body_font + 1)
        r.font.bold = True
        r.font.color.rgb = SHIPSY_NAVY

    def _bullets_text(left, top, width, height, items, font_pt):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        # Cap bullets so they fit
        for i, item in enumerate(items[:max_bullets]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.space_before = Pt(2)
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = "• " + item
            r.font.size = Pt(font_pt)
            r.font.color.rgb = DARK_TEXT

    # ----- Header row -----
    headers = [
        ("Business Process", MATRIX_PROC_BG),
        ("As Is",            MATRIX_ASIS_BG),
        ("To-Be",            MATRIX_TOBE_BG),
        ("Impact",           MATRIX_IMPACT_BG),
    ]
    for ci, (label, bg) in enumerate(headers):
        _cell_rect(col_xs[ci], grid_top, col_ws[ci], header_h, bg)
        _header_text(col_xs[ci], grid_top, col_ws[ci], header_h, label)

    # ----- Body rows -----
    for r_idx, row in enumerate(rows_data):
        y = grid_top + header_h + row_h * r_idx
        # Process column (light gray)
        _cell_rect(col_xs[0], y, col_ws[0], row_h, MATRIX_PROC_COL_BG)
        _proc_text(col_xs[0], y, col_ws[0], row_h, row.get("process", ""))
        # As-Is / To-Be / Impact (alternating row tint)
        body_bg = MATRIX_ROW_BG if r_idx % 2 == 1 else WHITE
        for ci, key in enumerate(["as_is", "to_be", "impact"], start=1):
            _cell_rect(col_xs[ci], y, col_ws[ci], row_h, body_bg)
            _bullets_text(col_xs[ci], y, col_ws[ci], row_h,
                          row.get(key, []), body_font)


def _add_phase_pill(slide, left, top, phase):
    color = {"Now": PHASE_NOW, "Phase 2": PHASE_PHASE2,
             "Future": PHASE_FUTURE}.get(phase, PHASE_FUTURE)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, Inches(0.85), Inches(0.24))
    _set_solid(pill, color)
    _add_textbox(slide, left, top, Inches(0.85), Inches(0.24),
                 phase, font_size=9, bold=True, color=WHITE)
    # The textbox doesn't centre in the pill — fix with paragraph alignment
    from pptx.enum.text import PP_ALIGN
    tf = pill.text_frame  # use the shape's own text_frame
    pill.text_frame.text = phase
    p = pill.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = WHITE
    return pill


def add_ai_use_cases_slide(prs, flow):
    """Slide: AI Use Cases · AgentFleet for {client}.

    A 2x3 grid of agent cards. Each card has agent name, tagline,
    client-specific scenario, and a phase pill (Now / Phase 2 / Future).
    """
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    sw, sh = prs.slide_width, prs.slide_height

    # Title bar
    title = f"AI Use Cases · AgentFleet for {flow['client']['name']}"
    subtitle = "Shipsy AI capabilities mapped to your specific operational scenarios"
    _add_title_bar(slide, sw, title, subtitle)

    cards = list(flow.get("ai_use_cases", []))[:6]
    if not cards:
        cards = [{
            "agent": "(none)", "tagline": "—",
            "scenario": "no ai_use_cases in flow JSON",
            "phase": "Future",
        }]

    # Grid: 3 columns × 2 rows
    grid_top = Inches(1.18)
    grid_left = Inches(0.4)
    grid_w = sw - Inches(0.8)
    grid_h = sh - grid_top - Inches(0.5)

    cols = 3
    rows = 2
    cell_w = (grid_w - Inches(0.30) * (cols - 1)) / cols
    cell_h = (grid_h - Inches(0.25) * (rows - 1)) / rows

    for i, card in enumerate(cards):
        c = i % cols
        r = i // cols
        if r >= rows:
            break
        x = grid_left + (cell_w + Inches(0.30)) * c
        y = grid_top + (cell_h + Inches(0.25)) * r

        # Card background
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x, y, cell_w, cell_h)
        _set_solid(bg, RGBColor(0xF6, 0xF8, 0xFB))
        bg.line.color.rgb = RGBColor(0xD0, 0xDA, 0xE5)

        # Top accent strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       x, y, cell_w, Inches(0.10))
        _set_solid(strip, SHIPSY_BLUE)

        pad_x = Inches(0.18)
        pill_w = Inches(0.78)
        pill_h = Inches(0.26)

        # Phase pill — TOP-RIGHT corner of card so it doesn't collide with
        # the scenario text below
        phase = card.get("phase", "Future")
        pill_color = {"Now": PHASE_NOW, "Phase 2": PHASE_PHASE2,
                      "Future": PHASE_FUTURE}.get(phase, PHASE_FUTURE)
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + cell_w - pill_w - Inches(0.18),
            y + Inches(0.20),
            pill_w, pill_h,
        )
        _set_solid(pill, pill_color)
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        ptf = pill.text_frame
        ptf.text = phase
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ptf.margin_left = Inches(0.04)
        ptf.margin_right = Inches(0.04)
        ptf.margin_top = Inches(0.02)
        ptf.margin_bottom = Inches(0.02)
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        for run in pp.runs:
            run.font.size = Pt(9)
            run.font.bold = True
            run.font.color.rgb = WHITE

        # Agent name — auto-shrink so long names fit on one line beside the
        # top-right phase pill. The column has ~1.6" of available width.
        agent_text = card["agent"]
        n_chars = len(agent_text)
        if n_chars <= 8:
            agent_font = 18
        elif n_chars <= 12:
            agent_font = 16
        elif n_chars <= 16:
            agent_font = 13
        else:
            agent_font = 11
        agent_w = cell_w - pad_x * 2 - pill_w - Inches(0.10)
        _add_textbox(slide, x + pad_x, y + Inches(0.20),
                     agent_w, Inches(0.34),
                     agent_text, font_size=agent_font, bold=True,
                     color=SHIPSY_NAVY)
        # Tagline
        _add_textbox(slide, x + pad_x, y + Inches(0.55),
                     cell_w - pad_x * 2, Inches(0.24),
                     card["tagline"], font_size=10, color=MUTED, bold=True)

        # Divider
        div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x + pad_x, y + Inches(0.82),
                                     cell_w - pad_x * 2, Inches(0.015))
        _set_solid(div, RGBColor(0xD9, 0xDF, 0xE6))

        # Scenario text — fills the rest of the card
        scenario_top = y + Inches(0.90)
        scenario_h = cell_h - Inches(0.90) - Inches(0.12)
        _add_textbox(slide, x + pad_x, scenario_top,
                     cell_w - pad_x * 2, scenario_h,
                     card["scenario"], font_size=9, color=DARK_TEXT)

    # Footer band
    foot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, sh - Inches(0.32), sw, Inches(0.32))
    _set_solid(foot, SHIPSY_CYAN)
    _add_textbox(slide, Inches(0.4), sh - Inches(0.30),
                 sw - Inches(0.8), Inches(0.28),
                 f"{flow['client']['name']} · "
                 f"{flow['client']['industry_vertical']} · "
                 "AI Use Cases · Shipsy",
                 font_size=10, color=WHITE, bold=True)

    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def render(flow_path: Path, out_path: Path) -> None:
    flow = json.loads(flow_path.read_text())

    # Start with a fresh widescreen 10" x 5.625" deck — no fixed template.
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Slide 1: dynamic Process Flow with Shipsy
    render_process_flow_v2(prs, flow)
    # Slide 2: Process Mapping & Capability Enhancement (matrix)
    add_process_mapping_slide(prs, flow)
    # Slide 3: AI Use Cases · AgentFleet
    add_ai_use_cases_slide(prs, flow)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    _strip_printer_settings(out_path)
    _libreoffice_roundtrip(out_path)
    print(f"Wrote {out_path}")


def _libreoffice_roundtrip(pptx_path: Path) -> None:
    """Round-trip the PPTX through LibreOffice so the OOXML produced by
    python-pptx (which Google Slides occasionally rejects) gets rewritten
    through a stricter serializer. Silent no-op if soffice isn't installed.
    """
    import shutil
    import subprocess
    import tempfile

    soffice = shutil.which('soffice') or shutil.which('libreoffice')
    if not soffice:
        return

    with tempfile.TemporaryDirectory() as td:
        try:
            subprocess.run(
                [soffice, '--headless', '--convert-to', 'pptx',
                 '--outdir', td, str(pptx_path)],
                check=True, capture_output=True, timeout=60,
            )
            roundtripped = Path(td) / pptx_path.name
            if roundtripped.exists() and roundtripped.stat().st_size > 0:
                shutil.move(str(roundtripped), str(pptx_path))
        except (subprocess.SubprocessError, FileNotFoundError):
            # Round-trip is best-effort; the strip-printer-settings step
            # above already handles the most common Google Slides blocker.
            return


def _strip_printer_settings(pptx_path: Path) -> None:
    """Remove printerSettings1.bin part — Google Slides rejects PPTX files
    that carry the legacy Windows-PowerPoint printer-config blob.
    Strips the part itself, the Default ContentType, and the relationship.
    """
    import zipfile
    import re
    import shutil

    tmp = pptx_path.with_suffix('.pptx.tmp')
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name in zin.namelist():
                if 'printerSettings' in name:
                    continue
                data = zin.read(name)
                if name == '[Content_Types].xml':
                    txt = data.decode('utf-8')
                    txt = re.sub(
                        r'<Default Extension="bin" ContentType="application/vnd\.openxmlformats-officedocument\.presentationml\.printerSettings"/>',
                        '', txt)
                    data = txt.encode('utf-8')
                elif name == 'ppt/_rels/presentation.xml.rels':
                    txt = data.decode('utf-8')
                    txt = re.sub(
                        r'<Relationship Id="[^"]+" Type="http://schemas\.openxmlformats\.org/officeDocument/2006/relationships/printerSettings" Target="[^"]+"/>',
                        '', txt)
                    data = txt.encode('utf-8')
                zout.writestr(name, data)
    shutil.move(tmp, pptx_path)


def main():
    if len(sys.argv) != 3:
        sys.stderr.write("usage: render.py <flow.json> <output.pptx>\n")
        sys.exit(2)
    render(Path(sys.argv[1]), Path(sys.argv[2]))


if __name__ == "__main__":
    main()
